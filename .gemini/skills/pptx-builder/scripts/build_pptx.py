import sys
import json
from pptx import Presentation
from pptx.util import Inches

def update_presentation(json_str):
    try:
        config = json.loads(json_str)
        template_path = config['template_path']
        output_path = config['output_path']
        updates = config.get('updates', [])

        prs = Presentation(template_path)

        for update in updates:
            slide_idx = update.get('slide_index') # 0-based index
            if slide_idx is None or slide_idx >= len(prs.slides):
                continue
            
            slide = prs.slides[slide_idx]

            # 1. Insert image (graph)
            img_path = update.get('image_path')
            if img_path:
                # Optional: Delete existing chart to make room for the new image
                shape_to_delete_name = update.get('replace_shape_name')
                if shape_to_delete_name:
                    for shape in list(slide.shapes):
                        if shape.name == shape_to_delete_name:
                            sp = shape._element
                            sp.getparent().remove(sp)

                # Default position: Left 1 inch, Top 1.5 inches, Width 8 inches
                # Can be overridden by config
                left = Inches(update.get('img_left_inches', 1))
                top = Inches(update.get('img_top_inches', 1.5))
                width = Inches(update.get('img_width_inches', 8))
                
                try:
                    slide.shapes.add_picture(img_path, left, top, width=width)
                except Exception as e:
                    print(f"Warning: Could not add image {img_path} to slide {slide_idx}. Error: {e}")

            # 2. Update 'Findings' text
            finding_text = update.get('finding_text')
            if finding_text:
                found = False
                # Try to find the existing 'Findings' text box and replace its content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and "Findings" in shape.text:
                        shape.text = f"Findings\n{finding_text}"
                        found = True
                        break
                
                # If no existing findings box is found, create a new one at the bottom
                if not found:
                    txBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
                    tf = txBox.text_frame
                    tf.text = f"Findings\n{finding_text}"

        prs.save(output_path)
        print(f"Success: Presentation saved to {output_path}")

    except Exception as e:
        print(f"Error updating presentation: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description='Update a PPTX template with graphs and text.')
    parser.add_argument('json_params', type=str, help='JSON string containing configuration.')
    args = parser.parse_args()
    
    update_presentation(args.json_params)
