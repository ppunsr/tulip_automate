import collections 
import collections.abc
from pptx import Presentation
import sys

def analyze_presentation(filepath):
    try:
        prs = Presentation(filepath)
        print(f"Analyzing: {filepath}")
        print(f"Total slides: {len(prs.slides)}\n")
        
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i + 1} ---")
            
            # Extract text
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            
            if texts:
                print("Text content found:")
                for t in texts:
                    # truncate long text for overview
                    preview = t.replace('\n', ' ')
                    if len(preview) > 100:
                        preview = preview[:97] + "..."
                    print(f"  - {preview}")
            else:
                print("No text content found.")
                
            # Look for placeholders/shapes for images
            print("Placeholders & Shapes:")
            for shape in slide.shapes:
                shape_type = shape.shape_type
                name = shape.name
                if shape.is_placeholder:
                    print(f"  - Placeholder: {name} (Type: {shape_type}, ID: {shape.placeholder_format.idx})")
                else:
                    print(f"  - Shape: {name} (Type: {shape_type})")
                    
            print("")
            
    except Exception as e:
        print(f"Error analyzing presentation: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        analyze_presentation(sys.argv[1])
    else:
        print("Please provide a path to a pptx file.")
